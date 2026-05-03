from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN


def create_portfolio():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(title, body_lines):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        body_shape = slide.placeholders[1]
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        tf = body_shape.text_frame
        tf.clear()
        for line in body_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.space_after = Pt(6)
        return slide

    # 1. 表纸
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "氏名：[英文氏名] / [日本語氏名]"
    slide.placeholders[1].text = "IoT / XR / Robotics　システムインテグレーションエンジニア\nEmail：[メールアドレス]\nWeb：[個人ウェブサイトURL（任意）]"
    # 2. 技术总括
    add_slide("技術総括", [
        "■ 技術領域",
        "・IoT / 組み込み　・XR / 3D可視化　・コンピュータビジョン / AI　・システムインテグレーション",
        "",
        "■ コア能力（強み）",
        "・システム設計　・マルチデバイス連携　・現場構築・導入　・トラブルシューティング / 最適化",
        "",
        "■ 技術スタック（主要抜粋）",
        "・[使用する主要技術スタックを簡潔に記載]"
    ])
    # 3. 项目1
    add_slide("プロジェクト1：スマート農業 IoT 自動化管理システム", [
        "【概要】農業環境モニタリングおよび自動制御システム。センサーにて温湿度等を収集し、遠隔監視と自動制御を実現。",
        "【担当】システム全体設計 / センサー・組み込み開発 / IoT通信・遠隔制御 / データ可視化UI / 現場導入・調整 / 保守運用",
        "【技術】IoT / センサーネットワーク、MCU、Python/C++、Web可視化",
        "【構成】センサー → MCU制御 → データ収集 → クラウド処理 → 可視化UI → 制御実行（灌漑/換気）",
        "【強み】マルチセンサー連携制御 / リアルタイムデータ収集・遠隔管理 / 自動制御ロジック / 長期安定稼働",
        "【実際の効果】※[現場写真 / 制御画面スクリーンショット / デバイス写真を挿入]",
        "【成果】農業環境の自動管理実現 / 管理効率向上・人件費削減 / 長期安定稼働を達成"
    ])
    # 4. 项目2
    add_slide("プロジェクト2：大型没入型インタラクティブ展示システム（XR / 5D）", [
        "【概要】展示館向け没入型システム。マルチスクリーン融合、リアルタイムレンダリング、インタラクティブ制御で高没入感を実現。",
        "【担当】マルチメディア設計 / マルチチャンネルプロジェクション開発 / Unity3D/OpenGL開発 / インタラクティブ開発 / 現場機器統合",
        "【技術】Unity3D / OpenGL、OpenCV、マルチプロジェクション、XRインタラクティブ",
        "【構成】入力（カメラ/センサー） → 認識処理 → レンダリング → マルチスクリーン出力 → フィードバック",
        "【強み】シームレスなプロジェクション融合 / リアルタイムフィードバック / 複雑な現場で高安定性 / 大規模システム統合力",
        "【実際の効果】※[展示館写真 / 投影画面 / システム現場写真を挿入]",
        "【成果】大規模施設へ導入成功 / 高没身体験の提供 / 複雑環境での安定稼働"
    ])
    # 5. 项目3
    add_slide("プロジェクト3：AIスマート除草ロボットシステム", [
        "【概要】コンピュータビジョンと自動制御による農業ロボット。画像認識で雑草を検出し自動処理、農業自動化を実現。",
        "【担当】システム設計 / YOLO物体認識適用 / ROS2制御開発 / 認識と制御ロジック統合 / プロトタイプ開発・テスト",
        "【技術】YOLO（物体検出）、OpenCV、ROS2（ロボット制御）、Python",
        "【構成】カメラ → 画像認識（AI） → 判断ロジック → 制御システム → 実行装置（除草）",
        "【強み】AI認識と自動制御の融合 / 農業作業の自動化 / スマート農機への拡張性",
        "【実際の効果】※[ロボット写真 / テスト風景 / 認識画面スクリーンショットを挿入]",
        "【成果】プロトタイプ開発完了 / 農業自動化でのAI活用実証 / 継続的最適化中"
    ])
    # 6. 总结
    add_slide("まとめ", [
        "■ 技術力の総括",
        "・IoT + XR + AIの融合　・システムインテグレーション能力　・現場導入・実装能力",
        "",
        "■ 強み",
        "・多分野にわたる開発経験　・設計から保守運用まで全工程対応　・実際のプロジェクト導入（落地）経験が豊富",
        "",
        "■ 今後の志向（農業分野）",
        "・IoT / AIを農業の自動化に応用したい　・生産効率の向上 / 知能化の推進に貢献する"
    ])
    prs.save('ポートフォリオ.pptx')
    print("PPT生成完了！ファイル名：ポートフォリオ.pptx")


create_portfolio()